from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import json
import re
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


#TODO Extract styles 


# Input content as a string
message = """
<!-- filename: Puppies_HELP_ME.docx --><head>{"bold": false, "italicised": false, "underline": false, "highlight": "none", "color": "black", "font": "Times New Roman", "fontSize": 12, "alignment": "justified"}</head><doc> <title>The Joyful World of Puppies</title>
<paragraph><span style="{"bold": true, "font": "Times New Roman", "fontSize": 12}">Introduction</span>: <sentence type="topic">The presence of puppies is synonymous with boundless energy, unconditional love, and unwavering companionship.</sentence> From their earliest moments, puppies captivate us with their innocent expressions and playful antics. This essay explores the development stages of puppies, their socialization needs, and the joy they bring into human lives.</paragraph>
 
<paragraph><span style="{"bold": false, "font": "Times New Roman", "fontSize": 12}">Embarking on Life</span>: <sentence type="topic">The journey of a puppy begins even before they open their eyes.</sentence> Born blind and deaf, puppies rely entirely on their mother for warmth and nutrition. During the first few weeks, rapid development takes place as they start to explore their surroundings, develop their senses, and establish bonds with their littermates and human caregivers.</paragraph>
 
<paragraph><span style="{"bold": false, "font": "Times New Roman", "fontSize": 12}">Growth and Socialization</span>: <sentence type="topic">As puppies grow, they undergo significant changes that prepare them for the world ahead.</sentence> Between three to twelve weeks of age, they experience a socialization period which is critical for their behavioral development. Positive interactions with people, other animals, and various environments during this phase are crucial for raising well-adjusted adult dogs. Moreover, this is the time when they begin basic training, learning commands such as "sit," "stay," and "come."</paragraph>
 
<paragraph><span style="{"bold": false, "font": "Times New Roman", "fontSize": 12}">The Bond with Humans</span>: <sentence type="topic">The bond formed between puppies and their owners is one of the most extraordinary aspects of the human-animal connection.</sentence> As pack animals, dogs have an innate ability to sense human emotions, often acting as emotional anchors. Raising a puppy offers both challenges and immense rewards, contributing to the emotional and physical well-being of their owners through companionship and love.</paragraph>
 
<paragraph><span style="{"bold": false, "font": "Times New Roman", "fontSize": 12}">Health and Care</span>: <sentence type="topic">Ensuring the health and happiness of a puppy requires dedication and knowledge.</sentence> This includes providing a balanced diet, regular exercise, and routine veterinary care. Vaccinations, spaying and neutering, and preventive measures against parasites are all part of responsible puppy care. Educating oneself about the specific needs of the puppy"s breed is also vital to offer the best quality of life.</paragraph>
 
<paragraph><span style="{"bold": true, "font": "Times New Roman", "fontSize": 12}">Conclusion</span>: <sentence type="conclusion">Puppies not only bring joy and liveliness into our homes but also teach us about responsibility, compassion, and unconditional love.</sentence> As they evolve from helpless bundles of fur to loyal and devoted companions, we, in turn, grow alongside them. The enriching experience of raising a puppy is one that shapes the lives of all who are fortunate enough to undertake it.</paragraph>
 
</doc>
"""

# Helper function to set run styles
def apply_styles(run, styles, default_styles):
    run.bold = styles.get("bold", default_styles["bold"])
    run.italic = styles.get("italicised", default_styles["italicised"])
    run.underline = styles.get("underline", default_styles["underline"])
    # Setting font properties. The default font and font size are extracted from the default_styles
    run.font.name = styles.get("font", default_styles["font"])
    run.font.size = Pt(styles.get("fontSize", default_styles["fontSize"]))
    # Add more style options here (font color, highlight, etc.) as needed
    
def docx_gen(input_content):
        # Extract the filename from the input content
    filename_match = re.search(r'<!-- filename: (.*?).docx -->', input_content)
    if not filename_match:
        raise ValueError("Filename not found in the input content.")

    filename = filename_match.group(1) + ".docx"
    header_match = re.search(r'<head>(.*?)<\/head>', input_content)

    if not header_match:
        raise ValueError("Header styles not found in the input content.")

    # Default styling from the header
    default_styles = json.loads(header_match.group(1))

    # Create a new Word document
    doc = Document()

    # Regex patterns for different components
    title_pattern = r'<title>(.*?)<\/title>'
    paragraph_pattern = r'<paragraph>(.*?)<\/paragraph>'
    span_pattern = r'<span style="(.*?)">(.*?)<\/span>'

    # Extract the title and add it to the document
    title_match = re.search(title_pattern, input_content)
    if title_match:
        doc.add_heading(title_match.group(1), 0)

    # Extract paragraphs and add them to the document
    paragraphs = re.findall(paragraph_pattern, input_content)
    for para in paragraphs:
        # First, deal with the <sentence type="..."> tags
        para = re.sub(r'<sentence type="[^"]+">(.*?)<\/sentence>', r'\1', para)
        
        span_matches = re.findall(span_pattern, para)
        p = doc.add_paragraph()
        for span_style, text in span_matches:
            span_style_dict = json.loads(span_style.replace("'", "\""))
            run = p.add_run(text)
            apply_styles(run, span_style_dict, default_styles)  # Now passing default_styles as well
            p.add_run(' ')  # For space between spans
        
        # Add rest of the paragraph text if there's any outside of span tags
        additional_text = re.sub(span_pattern, '', para)
        if additional_text.strip():
            p.add_run(additional_text.strip())

    # Save the document
    script_dir = os.path.dirname(os.path.abspath(__file__))
    filepath = f'../../files/docs/{filename}'

    full_path = os.path.join(script_dir, filepath)
    doc.save(filename)

def pptx_gen(input_content):
    filename_match = re.search(r'<!-- filename: (.*?).pptx -->', input_content)
    if not filename_match:
        raise ValueError("Filename not found in the input content.")

    filename = filename_match.group(1) + ".pptx"
    header_match = re.search(r'<head>(.*?)<\/head>', input_content)
    if not header_match:
        raise ValueError("Header styles not found in the input content.")
    header_styles = header_match.group(1).replace("'", '"')

    if not header_match:
        raise ValueError("Header styles not found in the input content.")
    
    default_styles = json.loads(header_match.group(1))
    
    presentation = Presentation()

    slide_pattern = r'<slide>(.*?)<\/slide>'
    title_pattern = r'<title>(.*?)<\/title>'
    bullet_pattern = r'<bullet>(.*?)<\/bullet>'

    # Extract slides from the input content
    slides = re.findall(slide_pattern, input_content, re.DOTALL)
    
    for slide_content in slides:
        # Add a new slide to the presentation
        slide_layout = presentation.slide_layouts[5]  # Choosing a blank slide layout
        slide = presentation.slides.add_slide(slide_layout)

        # Extract the title and add it to the slide
        title_match = re.search(title_pattern, slide_content)
        if title_match:
            title = title_match.group(1)
            title_placeholder = slide.shapes.title
            title_placeholder.text = title
            # Apply default styles to the title - we'll define a function to do this
        
        # Extract bullet points and add them to the slide
        bullet_points = re.findall(bullet_pattern, slide_content)
        if bullet_points:
            content_placeholder = slide.placeholders[1]  # Assuming placeholder 1 is the content box
            for bullet in bullet_points:
                p = content_placeholder.text_frame.add_paragraph()
                p.text = bullet
                # Apply default styles to the bullet points - we'll define a function to do this
    script_dir = os.path.dirname(os.path.abspath(__file__))
    filepath = f'../../files/docs/{filename}'

    full_path = os.path.join(script_dir, filepath)
    presentation.save(filename)
    

def gen_file(input_content):
    filename_match = re.search(r'<!-- filename: (.*?) -->', input_content)
    if not filename_match:
        raise ValueError("Filename not found in the input content.")
    
    filename = filename_match.group(1)
    if filename.endswith('.docx'):
        docx_gen(input_content)
    elif filename.endswith('.pptx'):
        pptx_gen(input_content)
    else:
        raise ValueError("Unknown document type.")
    

if __name__ == '__main__':
    gen_file(message)
